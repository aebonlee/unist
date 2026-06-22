#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
데이터 분석 실습 & 시각화 — 강의안 PowerPoint(.pptx) 생성기
public/slides.html 과 동일한 27장 구성·navy 디자인을 네이티브 PPTX로 재현한다.
출력: public/강의안_데이터분석_시각화.pptx
실행: python3 scripts/generate-pptx.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- 팔레트 (slides.html 과 동일) ----
NAVY    = RGBColor(0x1B, 0x2A, 0x4A)
NAVY3   = RGBColor(0x0D, 0x16, 0x26)
NAVY2   = RGBColor(0x14, 0x20, 0x36)
ACCENT  = RGBColor(0x3D, 0x7D, 0xFF)
ACCENTS = RGBColor(0x6F, 0xA0, 0xFF)
INK     = RGBColor(0xEE, 0xF2, 0xFB)
MUTED   = RGBColor(0x9F, 0xB0, 0xCF)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
CARD    = RGBColor(0x22, 0x30, 0x4E)
CARDLN  = RGBColor(0x3A, 0x4A, 0x6E)
LAB_BG  = ACCENT

FONT = "맑은 고딕"          # Windows/Office 기본 한글 폰트
FONT_MONO = "Consolas"

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ---------- helpers ----------
def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    bg.shadow.inherit = False
    # 우상단 은은한 톤 박스
    glow = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.3), Inches(-1), Inches(6), Inches(4))
    glow.fill.solid(); glow.fill.fore_color.rgb = RGBColor(0x24, 0x37, 0x60)
    glow.line.fill.background(); glow.shadow.inherit = False
    return s


def _set_font(run, size, color, bold=False, italic=False, mono=False):
    f = run.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.name = FONT_MONO if mono else FONT
    f.color.rgb = color
    # 한글 글꼴(EastAsian) 명시
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
    ea.set('typeface', FONT_MONO if mono else FONT)


def txt(s, l, t, w, h, runs, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, line_spacing=1.15, space_after=4):
    """runs: 문자열 또는 [(text, {opts})...] 리스트. 여러 문단은 '\n' 분리."""
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [(runs, {})]
    first = True
    for seg, opts in runs:
        lines = seg.split('\n')
        for i, ln in enumerate(lines):
            if first and i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.alignment = align
            p.line_spacing = line_spacing
            p.space_after = Pt(space_after)
            r = p.add_run(); r.text = ln
            _set_font(r, opts.get('size', size), opts.get('color', color),
                      opts.get('bold', bold), opts.get('italic', False), opts.get('mono', False))
        first = False
    return tb


def eyebrow(s, label):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.93), Inches(0.62),
                             Inches(0.05 + 0.105 * len(label) + 0.6), Inches(0.42))
    box.fill.solid(); box.fill.fore_color.rgb = CARD
    box.line.color.rgb = CARDLN; box.line.width = Pt(0.75); box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = False
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label.upper()
    _set_font(r, 11, ACCENTS, bold=True)


def header(s, label, title, sub=None):
    eyebrow(s, label)
    txt(s, Inches(0.9), Inches(1.18), Inches(11.5), Inches(1.0),
        title, size=33, color=WHITE, bold=True)
    if sub:
        txt(s, Inches(0.93), Inches(2.12), Inches(11.6), Inches(0.7),
            sub, size=15, color=MUTED)


def chrome(s, n, total):
    txt(s, Inches(0.9), Inches(7.02), Inches(3), Inches(0.4),
        [("Data", {'bold': True, 'color': INK, 'size': 12}),
         (" Lab", {'bold': True, 'color': ACCENTS, 'size': 12})])
    txt(s, Inches(10.5), Inches(7.02), Inches(2.0), Inches(0.4),
        f"{n} / {total}", size=12, color=MUTED, align=PP_ALIGN.RIGHT)
    # 진행바
    bar_w = EMU_W
    fill = int(EMU_W * n / total)
    base = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, bar_w, Pt(3))
    base.fill.solid(); base.fill.fore_color.rgb = NAVY2; base.line.fill.background(); base.shadow.inherit = False
    prog = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Emu(fill), Pt(3))
    prog.fill.solid(); prog.fill.fore_color.rgb = ACCENT; prog.line.fill.background(); prog.shadow.inherit = False


def card(s, l, t, w, h, title, body, num=None):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    c.fill.solid(); c.fill.fore_color.rgb = CARD
    c.line.color.rgb = CARDLN; c.line.width = Pt(0.75); c.shadow.inherit = False
    tf = c.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.18); tf.margin_bottom = Inches(0.14)
    p = tf.paragraphs[0]; p.line_spacing = 1.1
    if num:
        r = p.add_run(); r.text = num + "  "; _set_font(r, 15, ACCENTS, bold=True)
    r = p.add_run(); r.text = title; _set_font(r, 14, WHITE, bold=True)
    bp = tf.add_paragraph(); bp.line_spacing = 1.25; bp.space_before = Pt(6)
    r = bp.add_run(); r.text = body; _set_font(r, 11.5, MUTED)


def callout(s, l, t, w, segs, h=Inches(0.95)):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, Inches(0.06), h)
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background(); bar.shadow.inherit = False
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l + Inches(0.06), t, w - Inches(0.06), h)
    box.fill.solid(); box.fill.fore_color.rgb = CARD
    box.line.fill.background(); box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.22)
    p = tf.paragraphs[0]; p.line_spacing = 1.3
    for seg, o in segs:
        r = p.add_run(); r.text = seg
        _set_font(r, o.get('size', 13), o.get('color', INK), o.get('bold', False))


def bullets(s, l, t, w, h, items, size=14):
    """items: [(text, sub|None, [(part,bold)...]|None)] — 간단히 (text, sub)만 받음"""
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        main = it[0]; sub = it[1] if len(it) > 1 else None
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = 1.2; p.space_after = Pt(8)
        r = p.add_run(); r.text = "◆ "; _set_font(r, size - 2, ACCENT, bold=True)
        r = p.add_run(); r.text = main; _set_font(r, size, INK, bold=True)
        if sub:
            sp = tf.add_paragraph(); sp.line_spacing = 1.15; sp.space_after = Pt(8)
            r = sp.add_run(); r.text = "    " + sub; _set_font(r, size - 3, MUTED)


def table(s, l, t, w, headers, rows, col_w=None, fs=12, row_h=0.42):
    nrow = len(rows) + 1; ncol = len(headers)
    h = Inches(row_h * nrow)
    gt = s.shapes.add_table(nrow, ncol, l, t, w, h).table
    if col_w:
        total = sum(col_w)
        for ci, cw in enumerate(col_w):
            gt.columns[ci].width = Emu(int(w * cw / total))
    # 헤더
    for ci, hd in enumerate(headers):
        cell = gt.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY2
        cell.margin_left = Inches(0.1); cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]; r = p.add_run(); r.text = hd
        _set_font(r, fs, ACCENTS, bold=True)
    # 본문
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = gt.cell(ri, ci)
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY if ri % 2 else NAVY2
            cell.margin_left = Inches(0.1); cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]; p.line_spacing = 1.05
            r = p.add_run(); r.text = val
            _set_font(r, fs, WHITE if ci == 0 else MUTED, bold=(ci == 0))
    # 테두리 제거 스타일은 생략(기본 유지)
    return gt


def formula(s, l, t, w, text):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, Inches(0.62))
    box.fill.solid(); box.fill.fore_color.rgb = NAVY3
    box.line.color.rgb = CARDLN; box.line.width = Pt(0.75); box.shadow.inherit = False
    tf = box.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; _set_font(r, 18, ACCENTS, bold=True, mono=True)


def lab_tag(s, l, t):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, Inches(0.72), Inches(0.4))
    box.fill.solid(); box.fill.fore_color.rgb = LAB_BG; box.line.fill.background(); box.shadow.inherit = False
    tf = box.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "LAB"; _set_font(r, 13, WHITE, bold=True)


# =========================================================
# SLIDES
# =========================================================
TOTAL = 28
n = 0

# 1 · 표지
s = slide(); n += 1
txt(s, Inches(0.9), Inches(1.4), Inches(11), Inches(0.5),
    "2026 KERIS 교육과정(안) · 공통직무 · 중급", size=14, color=ACCENTS, bold=True)
txt(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.2),
    [("데이터 분석 실습\n& 시각화", {})], size=56, color=WHITE, bold=True, line_spacing=1.05)
txt(s, Inches(0.93), Inches(4.5), Inches(10), Inches(1.0),
    "엑셀 기반 EDA · 시각화 · 대시보드 구축 실습으로\n데이터 기반 의사결정 역량을 키우는 1일 6시간 과정",
    size=17, color=MUTED, line_spacing=1.3)
txt(s, Inches(0.93), Inches(6.6), Inches(12), Inches(0.5),
    [("강의안 · ", {'color': MUTED, 'size': 12}),
     ("Data Lab", {'color': WHITE, 'bold': True, 'size': 12}),
     (" · 드림아이티비즈(DreamIT Biz)  |  data.dreamitbiz.com", {'color': MUTED, 'size': 12})])

# 2 · 과정 개요
s = slide(); n += 1
header(s, "Course Overview", "과정 개요",
       "데이터를 정제하고 → 보여주고 → 결정으로 잇는 3단계 흐름을 하루에 익힙니다.")
cw = Inches(3.75); gap = Inches(0.2); top = Inches(2.85); ch = Inches(1.75)
card(s, Inches(0.9), top, cw, ch, "🎯 교육 목표",
     "데이터 분석·시각화 핵심 원리를 이해하고, 엑셀 실습으로 데이터 기반 의사결정 역량을 내재화합니다.")
card(s, Inches(0.9) + cw + gap, top, cw, ch, "⏱ 시간 / 레벨",
     "1일 6시간 · 중급\n모듈 3개 × 각 2.0H")
card(s, Inches(0.9) + (cw + gap) * 2, top, cw, ch, "🧰 사용 도구",
     "Microsoft Excel\n피벗테이블 · 조건부 서식 · 슬라이서")
callout(s, Inches(0.9), Inches(4.95), Inches(11.55),
        [("핵심 키워드 — ", {}), ("데이터 분석 · 엑셀 전처리 · EDA · 시각화 · 대시보드. ", {'bold': True, 'color': WHITE}),
         ("코딩 없이 엑셀만으로 분석 전 과정을 경험합니다.", {})])
chrome(s, n, TOTAL)

# 3 · 아젠다
s = slide(); n += 1
header(s, "Agenda", "오늘의 흐름", "탐색(EDA) → 표현(시각화) → 종합(대시보드)")
table(s, Inches(0.9), Inches(2.7), Inches(11.5),
      ["모듈", "주제", "핵심 역량", "시간"],
      [["모듈 1", "탐색적 데이터분석 (EDA)", "데이터 정제·구조 파악·요약통계·패턴 발굴", "2.0H"],
       ["모듈 2", "시각화의 원리 및 기획", "목적에 맞는 차트 선택과 정직한 표현", "2.0H"],
       ["모듈 3", "실무형 대시보드 구축", "KPI 종합 · 인터랙션 · 자동화", "2.0H"]],
      col_w=[1.3, 3.2, 5, 1], fs=13, row_h=0.62)
callout(s, Inches(0.9), Inches(5.3), Inches(11.55),
        [("분석의 끝은 ", {}), ("\"예쁜 차트\"가 아니라 \"더 나은 결정\"", {'bold': True, 'color': WHITE}),
         ("입니다. 세 모듈은 그 결정에 이르는 한 줄기 흐름입니다.", {})])
chrome(s, n, TOTAL)


# ---------- 디바이더 ----------
def divider(num, title_segs, sub, topics):
    global n
    s = slide(); n += 1
    big = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    big.fill.solid(); big.fill.fore_color.rgb = NAVY2; big.line.fill.background(); big.shadow.inherit = False
    txt(s, Inches(0.8), Inches(0.6), Inches(8), Inches(3),
        num, size=200, color=RGBColor(0x2A, 0x3C, 0x63), bold=True)
    txt(s, Inches(0.9), Inches(3.6), Inches(12), Inches(1.2),
        title_segs, size=44, color=WHITE, bold=True)
    txt(s, Inches(0.93), Inches(4.85), Inches(11.5), Inches(0.7),
        sub, size=17, color=MUTED)
    # topics chips
    x = Inches(0.9)
    for tp in topics:
        w = Inches(0.35 + 0.16 * len(tp))
        chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.7), w, Inches(0.46))
        chip.fill.solid(); chip.fill.fore_color.rgb = CARD
        chip.line.color.rgb = CARDLN; chip.line.width = Pt(0.75); chip.shadow.inherit = False
        tf = chip.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = tp; _set_font(r, 12, MUTED)
        x = Emu(int(x) + int(w) + Inches(0.18))
    chrome(s, n, TOTAL)


# ===== MODULE 1 =====
divider("01", [("탐색적 데이터분석 ", {'color': WHITE, 'bold': True}), ("(EDA)", {'color': ACCENTS, 'bold': True})],
        "데이터를 쪼개고 그려보며 패턴·이상치·가설을 발굴한다 · 2.0H",
        ["EDA의 본질", "문제·가설 수립", "지표 쪼개기", "시계열·상관·교차", "분포 확인"])

# M1-1 EDA란
s = slide(); n += 1
header(s, "Module 1", "EDA란 무엇인가",
       "통계학자 존 튜키 — 모델링 전에 데이터의 구조·패턴·이상치·관계를 탐색하는 과정")
bullets(s, Inches(0.9), Inches(2.7), Inches(5.6), Inches(3.5),
        [("데이터 품질 점검", "결측치·중복·잘못된 단위·비정상 값 발견"),
         ("가설 발굴", "‘이 두 변수는 관련 있을까?’ 새 질문 생성"),
         ("이상치 탐지", "평균을 왜곡하는 극단값·입력 오류"),
         ("분석 방향 설정", "어떤 변수를 깊게, 어떤 모델이 적합할지")])
table(s, Inches(6.8), Inches(2.7), Inches(5.6),
      ["구분", "EDA (탐색)", "CDA (확증)"],
      [["목적", "패턴·가설 발굴", "가설 검증"],
       ["가설", "분석 중 생성", "분석 전 수립"],
       ["도구", "시각화·요약통계", "가설검정·p-value"],
       ["비유", "탐정 현장조사", "법정 증거검증"]],
      col_w=[1.1, 2.2, 2.3], fs=11.5, row_h=0.5)
callout(s, Inches(0.9), Inches(6.1), Inches(11.55),
        [("실무에서 EDA는 분석 시간의 ", {}), ("60~80%", {'bold': True, 'color': WHITE}),
         ("를 차지합니다. ‘Garbage In, Garbage Out’을 막는 ", {}), ("첫 방어선", {'bold': True, 'color': WHITE}), (".", {})])
chrome(s, n, TOTAL)

# M1-2 문제·가설
s = slide(); n += 1
header(s, "Module 1", "문제 정의와 가설 수립",
       "분석은 ‘데이터’가 아니라 ‘문제’에서 출발한다 — 막연한 고민을 답할 수 있는 질문으로")
table(s, Inches(0.9), Inches(2.7), Inches(11.5),
      ["막연한 고민", "분석 가능한 질문"],
      [["\"매출이 왜 줄었지?\"", "최근 3개월간 어느 제품군·지역·채널에서 매출이 가장 감소했는가?"],
       ["\"고객이 이탈하는 것 같아\"", "가입 후 30일 내 재방문 안 한 고객의 비율과 특성은?"],
       ["\"마케팅이 효과가 있나?\"", "캠페인 전후 2주간 신규 전환율이 유의하게 달라졌는가?"]],
      col_w=[3.5, 8], fs=13, row_h=0.6)
callout(s, Inches(0.9), Inches(5.2), Inches(11.55),
        [("핵심은 ", {}), ("측정 가능한 지표 · 명확한 기간 · 비교 대상", {'bold': True, 'color': WHITE}),
         (". 좋은 가설 = 검증 가능 · 구체적 · 반증 가능 · 실행 연결.", {})])
chrome(s, n, TOTAL)

# M1-3 5단계
s = slide(); n += 1
header(s, "Module 1", "논리적 검증 프로세스", "문제 정의에서 결론까지 5단계")
formula(s, Inches(0.9), Inches(2.6), Inches(11.5),
        "문제 정의 → 가설 수립 → 지표 정의 → 데이터 검증 → 결론·액션")
table(s, Inches(0.9), Inches(3.6), Inches(11.5),
      ["단계", "핵심 질문", "산출물"],
      [["1 · 문제 정의", "무엇이 문제인가?", "분석 가능한 질문"],
       ["2 · 가설 수립", "왜 그럴 것 같은가?", "검증 가능한 가설 목록"],
       ["3 · 지표 정의", "무엇을 측정할 것인가?", "KPI · 측정 기준"],
       ["4 · 데이터 검증", "데이터가 가설을 지지하는가?", "EDA 결과 · 차트"],
       ["5 · 결론·액션", "그래서 무엇을 할 것인가?", "의사결정 · 권고안"]],
      col_w=[2.5, 5, 4], fs=12.5, row_h=0.52)
chrome(s, n, TOTAL)

# M1-4 드릴다운
s = slide(); n += 1
header(s, "Module 1", "지표 쪼개기 (Drill-down)",
       "‘무슨 일이 일어났는가’를 넘어 ‘왜 그랬는가’로 — 큰 숫자를 구성 요소로 분해")
formula(s, Inches(0.9), Inches(2.55), Inches(6), "매출 = 방문수 × 전환율 × 객단가")
table(s, Inches(0.9), Inches(3.45), Inches(6),
      ["분해 대상", "공식"],
      [["매출", "방문수 × 전환율 × 객단가"],
       ["활성 사용자", "신규 + 유지 + 복귀"],
       ["이익", "매출 − (고정비 + 변동비)"]],
      col_w=[2, 4], fs=12, row_h=0.5)
bullets(s, Inches(7.2), Inches(2.6), Inches(5.2), Inches(3),
        [("차원으로 쪼개기", "시간·고객·상품·채널/지역 관점"),
         ("다차원 교차", "‘지역 × 연령대’로 숨은 신호 포착"),
         ("심슨의 역설 주의", "전체는 평탄해도 부분이 상쇄될 수 있음")])
callout(s, Inches(0.9), Inches(6.0), Inches(11.55),
        [("엑셀 ", {}), ("피벗테이블", {'bold': True, 'color': WHITE}),
         (" — 지표를 값, 1차 차원을 행, 2차 차원을 열에 두면 교차표 완성. 슬라이서로 필터링.", {})])
chrome(s, n, TOTAL)

# M1-5 LAB 시계열·집단
s = slide(); n += 1
eyebrow(s, "Module 1 · 실습")
lab_tag(s, Inches(0.9), Inches(1.2))
txt(s, Inches(1.75), Inches(1.18), Inches(11), Inches(0.9),
    "시계열 추세 · 집단별 비교", size=31, color=WHITE, bold=True)
txt(s, Inches(0.93), Inches(2.15), Inches(11.5), Inches(0.6),
    "시간에 따른 추세와 집단 간 차이를 엑셀로 직접 분석", size=15, color=MUTED)
cw = Inches(3.75); top = Inches(2.95); ch = Inches(2.05)
card(s, Inches(0.9), top, cw, ch, "시계열 추세",
     "꺾은선 차트 + 7일 이동평균(=AVERAGE)으로 단기 변동 속 장기 추세·계절성 파악", num="1")
card(s, Inches(0.9) + cw + gap, top, cw, ch, "요일별 특성",
     "피벗(행=요일, 값=매출 평균) → 막대그래프로 주말 vs 평일 비교", num="2")
card(s, Inches(0.9) + (cw + gap) * 2, top, cw, ch, "집단별 비교",
     "행=고객유형/지점, 합계와 평균을 함께 → 규모와 효율 동시에", num="3")
callout(s, Inches(0.9), Inches(5.3), Inches(11.55),
        [("그룹 비교 시 ", {}), ("합계만 보지 말고 평균·비율도 함께", {'bold': True, 'color': WHITE}),
         (". 규모가 커서 합계만 높고 효율(1인당)은 낮을 수 있습니다.", {})])
chrome(s, n, TOTAL)

# M1-6 LAB 상관·교차
s = slide(); n += 1
eyebrow(s, "Module 1 · 실습")
lab_tag(s, Inches(0.9), Inches(1.2))
txt(s, Inches(1.75), Inches(1.18), Inches(11), Inches(0.9), "상관분석 · 교차분석", size=31, color=WHITE, bold=True)
txt(s, Inches(0.93), Inches(2.15), Inches(11.5), Inches(0.6),
    "두 변수가 함께 움직이는가(상관) · 범주 간 관계가 있는가(교차)", size=15, color=MUTED)
table(s, Inches(0.9), Inches(2.95), Inches(5.6),
      ["상관계수 r", "해석"],
      [["0.7 ~ 1.0", "강한 양의 상관"],
       ["0.4 ~ 0.7", "뚜렷한 양의 상관"],
       ["−0.2 ~ 0.2", "거의 상관 없음"],
       ["−1.0 ~ −0.7", "강한 음의 상관"]],
      col_w=[2.3, 3.3], fs=12, row_h=0.5)
formula(s, Inches(0.9), Inches(5.2), Inches(5.6), "=CORREL(B2:B100, C2:C100)")
callout(s, Inches(6.9), Inches(2.95), Inches(5.5),
        [("⚠ 상관 ≠ 인과. ", {'bold': True, 'color': WHITE}),
         ("‘아이스크림 판매 ↔ 익사’는 강한 상관이지만 숨은 변수 ‘여름철 더위’가 둘을 끌어올린 것. r은 직선 관계만 측정하므로 반드시 산점도를 함께 그리세요.", {})],
        h=Inches(2.7))
chrome(s, n, TOTAL)

# M1-7 LAB 분포
s = slide(); n += 1
eyebrow(s, "Module 1 · 실습")
lab_tag(s, Inches(0.9), Inches(1.2))
txt(s, Inches(1.75), Inches(1.18), Inches(11.2), Inches(0.9),
    "분포 확인 — 히스토그램 · 박스플롯", size=29, color=WHITE, bold=True)
txt(s, Inches(0.93), Inches(2.15), Inches(11.5), Inches(0.6),
    "평균만 보면 데이터의 진짜 모습을 놓친다 — 분포의 형태를 직접 본다", size=15, color=MUTED)
bullets(s, Inches(0.9), Inches(2.95), Inches(5.7), Inches(3),
        [("평균의 함정", "9명 3천만+1명 30억 → 평균 3억, 누구도 안 받음. 중앙값을"),
         ("히스토그램", "대칭/치우침(왜도)/봉우리 수 → 분포 형태"),
         ("박스플롯", "사분위수·IQR·수염 밖의 점=이상치")])
table(s, Inches(6.9), Inches(2.95), Inches(5.5),
      ["분포 형태", "대표값"],
      [["종 모양(정규)", "평균"],
       ["오른쪽 꼬리(소득)", "중앙값"],
       ["이봉(두 집단)", "집단 분리 후"]],
      col_w=[3, 2.5], fs=12, row_h=0.5)
callout(s, Inches(6.9), Inches(5.0), Inches(5.5),
        [("앤스컴 콰르텟", {'bold': True, 'color': WHITE}),
         (" — 평균·분산·상관이 같아도 모양은 전혀 다름. 요약 전에 반드시 그려보라.", {})], h=Inches(1.3))
chrome(s, n, TOTAL)

# M1-8 체크리스트
s = slide(); n += 1
header(s, "Module 1 · 정리", "EDA 실무 체크리스트",
       "좋은 EDA는 차트를 많이 그린 것이 아니라, 올바른 질문을 던진 것")
bullets(s, Inches(0.9), Inches(2.75), Inches(5.7), Inches(3.5),
        [("분석 목적과 가설을 한 문장으로 적었는가?",),
         ("행/열 수·결측치·중복·단위를 점검했는가?",),
         ("핵심 지표를 구성 요소·차원으로 쪼갰는가?",),
         ("시간에 따른 추세·계절성을 확인했는가?",)])
bullets(s, Inches(6.8), Inches(2.75), Inches(5.7), Inches(3.5),
        [("분포(히스토그램·박스플롯)를 그렸는가?",),
         ("이상치가 오류인지 진짜 특이값인지 판단했는가?",),
         ("상관을 인과로 단정하지 않았는가?",),
         ("발견한 패턴을 새 질문으로 정리했는가?",)])
chrome(s, n, TOTAL)


# ===== MODULE 2 =====
divider("02", [("시각화의 원리 ", {'color': WHITE, 'bold': True}), ("및 기획", {'color': ACCENTS, 'bold': True})],
        "데이터가 담은 메시지를 가장 빠르고 정확하게 전달한다 · 2.0H",
        ["메시지·타겟", "차트 선택", "Data-Ink Ratio", "왜곡 회피", "윤리적 표현"])

# M2-1 목적
s = slide(); n += 1
header(s, "Module 2", "시각화의 목적 — 메시지와 타겟",
       "예쁜 그림이 아니라 결론을 보여주는 차트. 그리기 전에 ‘무엇을 / 누구에게’부터")
callout(s, Inches(0.9), Inches(2.6), Inches(11.55),
        [("차트 제목에 ‘월별 매출’이 아니라 ", {}), ("‘3분기 매출이 전년 대비 18% 감소했다’", {'bold': True, 'color': WHITE}),
         ("처럼 결론(So-What)을 쓰는 순간 불필요한 요소가 정리됩니다.", {})], h=Inches(0.9))
table(s, Inches(0.9), Inches(3.85), Inches(11.5),
      ["타겟", "관심사", "시각화 전략"],
      [["경영진", "결론·의사결정", "KPI 1~2개, 큰 숫자, 추세 한눈에. 디테일 제거"],
       ["실무자", "원인·세부", "차원별 분해, 비교 가능한 표·차트, 필터 제공"],
       ["외부/고객", "신뢰·맥락", "출처 명시, 단순·오해 없는 표현, 스토리 흐름"]],
      col_w=[1.8, 2.5, 7.2], fs=12.5, row_h=0.55)
chrome(s, n, TOTAL)

# M2-2 차트 선택
s = slide(); n += 1
header(s, "Module 2", "차트 선택 가이드",
       "‘내가 보여주려는 관계가 무엇인가’를 정하면 차트는 거의 자동으로 결정된다")
table(s, Inches(0.9), Inches(2.7), Inches(11.5),
      ["분석 목적", "권장 차트", "예시"],
      [["비교 (크기)", "막대그래프", "지점별 매출 비교"],
       ["추이 (시간)", "꺾은선그래프", "월별 방문자 추세"],
       ["구성 (비율)", "누적막대 · (제한적)파이", "채널별 매출 비중"],
       ["관계 (상관)", "산점도", "광고비 vs 매출"],
       ["분포 (퍼짐)", "히스토그램 · 박스플롯", "고객 연령 분포"]],
      col_w=[3, 4.5, 4], fs=12.5, row_h=0.5)
callout(s, Inches(0.9), Inches(5.7), Inches(11.55),
        [("규칙 — ", {'bold': True, 'color': WHITE}),
         ("비교는 막대, 추이는 선, 관계는 산점도, 분포는 히스토그램/박스플롯. 파이 남용·스파게티 차트를 피하세요.", {})])
chrome(s, n, TOTAL)

# M2-3 LAB 가독성
s = slide(); n += 1
eyebrow(s, "Module 2 · 실습")
lab_tag(s, Inches(0.9), Inches(1.2))
txt(s, Inches(1.75), Inches(1.18), Inches(11.2), Inches(0.9),
    "가독성 — Data-Ink Ratio", size=30, color=WHITE, bold=True)
txt(s, Inches(0.93), Inches(2.15), Inches(11.5), Inches(0.6),
    "좋은 차트는 더하는 게 아니라 덜어낸 결과. 잉크의 대부분을 데이터에만", size=15, color=MUTED)
bullets(s, Inches(0.9), Inches(2.95), Inches(5.7), Inches(3),
        [("불필요한 요소 제거", "격자선·테두리·배경·3D·범례·과도한 소수점"),
         ("강조색은 1개", "나머지 회색. 빨강-초록 조합 회피(색맹)"),
         ("시선 유도", "값 순 정렬·직접 라벨·제목에 결론")])
table(s, Inches(6.9), Inches(2.95), Inches(5.5),
      ["항목", "Before", "After"],
      [["격자선", "진한 검정", "옅은 회색"],
       ["색상", "무지개 7색", "강조 1색+회색"],
       ["범례", "별도 박스", "직접 라벨"],
       ["제목", "‘월별 매출’", "‘18% 감소’"]],
      col_w=[1.4, 2.1, 2], fs=11, row_h=0.48)
callout(s, Inches(0.9), Inches(5.95), Inches(11.55),
        [("‘더 이상 뺄 것이 없을 때 완성된다.’ 차트도 디자인도 마찬가지입니다.", {})])
chrome(s, n, TOTAL)

# M2-4 LAB 왜곡
s = slide(); n += 1
eyebrow(s, "Module 2 · 실습")
lab_tag(s, Inches(0.9), Inches(1.2))
txt(s, Inches(1.75), Inches(1.18), Inches(11.2), Inches(0.9),
    "정보 왜곡 오류 피하기", size=30, color=WHITE, bold=True)
txt(s, Inches(0.93), Inches(2.15), Inches(11.5), Inches(0.6),
    "가장 위험한 건 못생긴 차트가 아니라 사실을 왜곡하는 차트", size=15, color=MUTED)
table(s, Inches(0.9), Inches(2.95), Inches(11.5),
      ["왜곡 유형", "무엇이 문제인가", "올바른 대안"],
      [["y축 0 미시작", "작은 차이를 거대하게 과장", "막대는 반드시 0부터"],
       ["이중 축 오용", "스케일 조정해 가짜 상관 연출", "차트 분리 / 지수화"],
       ["3D 차트", "원근·기울기로 크기 왜곡", "평면 2D 사용"],
       ["면적/버블 과장", "지름→면적이 제곱으로 과장", "면적을 값에 비례"]],
      col_w=[2.5, 5, 4], fs=12, row_h=0.5)
callout(s, Inches(0.9), Inches(5.5), Inches(11.55),
        [("황금률 — ", {'bold': True, 'color': WHITE}),
         ("막대그래프 y축은 반드시 0에서 시작. 꺾은선은 아니어도 되지만 축 범위를 반드시 명시.", {})])
chrome(s, n, TOTAL)

# M2-5 LAB 윤리
s = slide(); n += 1
eyebrow(s, "Module 2 · 실습")
lab_tag(s, Inches(0.9), Inches(1.2))
txt(s, Inches(1.75), Inches(1.18), Inches(11.2), Inches(0.9),
    "윤리적이고 정확한 표현", size=30, color=WHITE, bold=True)
txt(s, Inches(0.93), Inches(2.15), Inches(11.5), Inches(0.6),
    "시각화의 마지막 책임은 정직함 — 맥락을 숨기면 거짓말과 같다", size=15, color=MUTED)
bullets(s, Inches(0.9), Inches(2.95), Inches(5.7), Inches(3),
        [("체리피킹 회피", "유리한 기간·집단만 고르지 않기"),
         ("기준선 조작 회피", "전년? 목표? 경쟁사? 기준 명시"),
         ("표본 왜곡 회피", "‘90% 만족’(n=10) → 표본 n 표기")])
card(s, Inches(6.9), Inches(2.95), Inches(5.5), Inches(2.3), "객관적 팩트 전달 3원칙",
     "1. 출처 명시 — 데이터 출처·수집 방법\n2. 기간·표본 명시 — 언제, 몇 건\n3. 단위·정의 명시 — 매출은 부가세 포함?")
callout(s, Inches(0.9), Inches(6.0), Inches(11.55),
        [("좋은 분석가는 ‘내 주장을 돕는 차트’가 아니라 ", {}),
         ("‘사실을 가장 정확히 보여주는 차트’", {'bold': True, 'color': WHITE}), ("를 만듭니다.", {})])
chrome(s, n, TOTAL)

# M2-6 치트시트
s = slide(); n += 1
header(s, "Module 2 · 정리", "‘이 데이터엔 이 차트’ 치트시트",
       "시각화는 메시지 전달 수단 — 목적이 차트를 결정한다")
table(s, Inches(0.9), Inches(2.7), Inches(11.5),
      ["보여주려는 것", "첫 번째 선택", "피할 것"],
      [["항목 간 크기 비교", "가로 막대(정렬)", "파이(5개 이상)"],
       ["시간에 따른 변화", "꺾은선", "막대(너무 많을 때)"],
       ["부분과 전체 비율", "누적막대 100%", "3D 파이"],
       ["두 변수 관계", "산점도", "이중 축"],
       ["값의 분포", "히스토그램/박스플롯", "평균만 표기"],
       ["단일 핵심 지표", "큰 숫자 + 스파크라인", "게이지/도넛 남용"]],
      col_w=[4, 4, 3.5], fs=12.5, row_h=0.5)
chrome(s, n, TOTAL)


# ===== MODULE 3 =====
divider("03", [("실무형 ", {'color': WHITE, 'bold': True}), ("대시보드 구축", {'color': ACCENTS, 'bold': True})],
        "한 화면에서 KPI를 모니터링하는 인터랙티브 대시보드 · 2.0H",
        ["인-셀 시각화", "KPI·레이아웃", "피벗·슬라이서", "자동 갱신", "종합"])

# M3-1 인-셀
s = slide(); n += 1
header(s, "Module 3", "표를 직관적으로 — 인-셀 시각화",
       "차트 없이 셀 안에서 값의 크기·추세·상태를 표현한다")
cw = Inches(3.75); top = Inches(2.85); ch = Inches(2.1)
card(s, Inches(0.9), top, cw, ch, "🎨 조건부 서식",
     "색조(값 높낮이) · 아이콘 집합(신호등·화살표) · 규칙(상위 10%·중복 강조)")
card(s, Inches(0.9) + cw + gap, top, cw, ch, "📊 데이터 막대",
     "셀 안에 막대를 그려 상대 크기 표시. ‘막대만 표시’로 숫자를 숨기면 깔끔")
card(s, Inches(0.9) + (cw + gap) * 2, top, cw, ch, "📈 스파크라인",
     "한 셀 미니 차트 — 꺾은선(추세)·열(크기)·승패(달성/미달)")
callout(s, Inches(0.9), Inches(5.25), Inches(11.55),
        [("목적은 ", {}), ("‘값을 읽지 않고도 패턴이 보이게’", {'bold': True, 'color': WHITE}),
         (" 하는 것. 색·아이콘은 2~3종 이내로 절제하세요.", {})])
chrome(s, n, TOTAL)

# M3-2 기획
s = slide(); n += 1
header(s, "Module 3", "대시보드 기획 — KPI와 레이아웃",
       "잘 그린 차트의 모음이 아니라, 의사결정을 돕는 한 장의 화면")
cw2 = Inches(3.7); top = Inches(2.8); ch = Inches(1.5)
card(s, Inches(0.9), top, cw2, ch, "한 화면 원칙", "스크롤 없이 핵심을 한 화면에", num="1")
card(s, Inches(0.9), top + ch + Inches(0.15), cw2, ch, "핵심 KPI 우선", "중요 지표를 상단·좌측에 크게", num="2")
card(s, Inches(0.9), top + (ch + Inches(0.15)) * 2, cw2, ch, "시선 흐름", "눈은 F자/Z자 — 좌상→우상→하단", num="3")
table(s, Inches(5.0), Inches(2.8), Inches(7.4),
      ["부서", "핵심 KPI"],
      [["영업", "월 매출액, 목표 달성률"],
       ["R&D", "프로젝트 진척률, 일정 준수율"],
       ["인사", "채용 충원율, 이직률"]],
      col_w=[2, 5.4], fs=13, row_h=0.55)
callout(s, Inches(5.0), Inches(4.9), Inches(7.4),
        [("KPI는 ‘많이 보여주기’가 아니라 ", {}), ("‘골라 보여주기’", {'bold': True, 'color': WHITE}),
         (". 10개 넘으면 아무것도 강조 안 된 것.", {})], h=Inches(1.1))
chrome(s, n, TOTAL)

# M3-3 와이어프레임
s = slide(); n += 1
header(s, "Module 3", "와이어프레임 설계 단계", "셀에 그리기 전에 종이/슬라이드에 박스로 먼저")
steps = [("목적 정의", "‘누가, 무엇을, 왜 보는가’를 한 문장으로"),
         ("KPI 선정", "핵심 3~5개로 압축 (많을수록 나쁨)"),
         ("영역 분할", "상단 KPI 카드 / 중앙 추세 차트 / 하단 상세 표"),
         ("종이 스케치", "박스 배치를 먼저 그려본다"),
         ("그리드 설정", "눈금선 끄고 셀 너비 정렬해 격자 레이아웃")]
y = Inches(2.75)
for i, (t, d) in enumerate(steps, 1):
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), y, Inches(0.5), Inches(0.5))
    circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT; circ.line.fill.background(); circ.shadow.inherit = False
    tf = circ.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(i); _set_font(r, 14, WHITE, bold=True)
    txt(s, Inches(1.65), y - Inches(0.02), Inches(10.5), Inches(0.6),
        [(t + "  ", {'bold': True, 'color': WHITE, 'size': 16}), ("— " + d, {'color': MUTED, 'size': 14})],
        anchor=MSO_ANCHOR.MIDDLE)
    y = Emu(int(y) + int(Inches(0.82)))
chrome(s, n, TOTAL)

# M3-4 LAB 인터랙션
s = slide(); n += 1
eyebrow(s, "Module 3 · 실습")
lab_tag(s, Inches(0.9), Inches(1.2))
txt(s, Inches(1.75), Inches(1.18), Inches(11.4), Inches(0.9),
    "피벗차트 · 슬라이서 · 타임라인", size=28, color=WHITE, bold=True)
txt(s, Inches(0.93), Inches(2.15), Inches(11.5), Inches(0.6),
    "피벗테이블 + 슬라이서를 연결하면 클릭만으로 여러 차트가 동시에 바뀐다", size=15, color=MUTED)
card(s, Inches(0.9), Inches(2.95), Inches(5.7), Inches(1.9), "🔗 슬라이서 (Slicer)",
     "범주(지역·채널·부서) 버튼 필터. 여러 피벗에 ‘보고서 연결’하면 한 클릭으로 전체 갱신")
card(s, Inches(6.75), Inches(2.95), Inches(5.65), Inches(1.9), "📅 시간 표시 막대 (Timeline)",
     "날짜 전용 필터 — 연/분기/월/일 단위로 기간을 슬라이드")
callout(s, Inches(0.9), Inches(5.2), Inches(11.55),
        [("핵심 — ", {'bold': True, 'color': WHITE}),
         ("슬라이서/타임라인을 모든 피벗 요소에 연결하면 클릭 한 번으로 대시보드 전체를 탐색하는 인터랙티브 화면이 됩니다.", {})])
chrome(s, n, TOTAL)

# M3-5 LAB 자동갱신
s = slide(); n += 1
eyebrow(s, "Module 3 · 실습")
lab_tag(s, Inches(0.9), Inches(1.2))
txt(s, Inches(1.75), Inches(1.18), Inches(11.2), Inches(0.9),
    "자동 갱신 리포트", size=30, color=WHITE, bold=True)
txt(s, Inches(0.93), Inches(2.15), Inches(11.5), Inches(0.6),
    "만들고 계속 쓰는 템플릿 — 수작업 업데이트는 곧 방치로 이어진다", size=15, color=MUTED)
bullets(s, Inches(0.9), Inches(2.95), Inches(6), Inches(3),
        [("표(Table)로 변환", "Ctrl+T — 데이터가 늘어도 범위 자동 확장"),
         ("모두 새로 고침", "원본만 갱신하면 피벗·차트 한 번에 업데이트"),
         ("GETPIVOTDATA", "피벗 값을 KPI 카드로 안전하게 끌어오기")])
callout(s, Inches(7.0), Inches(2.95), Inches(5.4),
        [("원본 갱신이 ", {}), ("‘모두 새로 고침’ 한 번", {'bold': True, 'color': WHITE}),
         ("으로 반영되도록 처음부터 설계하면, 매달 손대지 않아도 살아있는 리포트가 됩니다.", {})], h=Inches(2.2))
chrome(s, n, TOTAL)

# M3-6 종합 점검
s = slide(); n += 1
eyebrow(s, "Module 3 · 종합")
lab_tag(s, Inches(0.9), Inches(1.2))
txt(s, Inches(1.75), Inches(1.18), Inches(11.4), Inches(0.9),
    "부서 KPI 대시보드 점검", size=29, color=WHITE, bold=True)
txt(s, Inches(0.93), Inches(2.15), Inches(11.5), Inches(0.6),
    "인-셀 시각화 · 기획 · 인터랙션 · 자동화를 한 화면에", size=15, color=MUTED)
bullets(s, Inches(0.9), Inches(2.95), Inches(5.7), Inches(3.2),
        [("핵심 KPI가 상단·좌측에 우선 배치?",),
         ("스크롤 없이 한 화면에 들어오는가?",),
         ("슬라이서/타임라인이 모든 요소에 연결?",)])
bullets(s, Inches(6.8), Inches(2.95), Inches(5.7), Inches(3.2),
        [("원본 갱신이 ‘모두 새로 고침’만으로 반영?",),
         ("색·아이콘이 절제되어 강조가 분명?",),
         ("제목·기준 기간·갱신일자가 표시?",)])
chrome(s, n, TOTAL)


# ===== 회고 =====
s = slide(); n += 1
header(s, "Course Recap", "전체 과정 회고", "데이터를 다루는 세 단계를 차례로 익혔습니다")
table(s, Inches(0.9), Inches(2.7), Inches(11.5),
      ["모듈", "주제", "핵심 역량"],
      [["1", "탐색적 데이터분석 (EDA)", "데이터 정제 · 구조 파악 · 요약통계"],
       ["2", "시각화", "목적에 맞는 차트 선택과 정직한 표현"],
       ["3", "대시보드", "KPI 종합 · 인터랙션 · 자동화"]],
      col_w=[1, 4, 6.5], fs=13, row_h=0.6)
callout(s, Inches(0.9), Inches(5.2), Inches(11.55),
        [("현업 적용 — ", {'bold': True, 'color': WHITE}),
         ("① 작게 시작(KPI 3개) ② 자동 갱신을 기본으로 ③ 보는 사람 기준(1분 안에) ④ 분기마다 KPI 점검.", {})])
chrome(s, n, TOTAL)

# ===== 클로징 =====
s = slide(); n += 1
txt(s, Inches(0.9), Inches(1.4), Inches(11), Inches(0.5), "THANK YOU", size=14, color=ACCENTS, bold=True)
txt(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(2.0),
    "데이터 분석의 끝은\n더 나은 결정입니다", size=46, color=WHITE, bold=True, line_spacing=1.1)
txt(s, Inches(0.93), Inches(4.5), Inches(11), Inches(1.2),
    "대시보드는 그 결정을 빠르게 돕는 도구일 뿐.\n오늘 익힌 EDA → 시각화 → 대시보드를 현업의 한 장으로 옮겨보세요. 수고하셨습니다.",
    size=16, color=MUTED, line_spacing=1.35)
txt(s, Inches(0.93), Inches(6.5), Inches(12), Inches(0.5),
    [("data.dreamitbiz.com", {'color': ACCENTS, 'bold': True, 'size': 13}),
     ("   ·   드림아이티비즈(DreamIT Biz)   ·   aebon@dreamitbiz.com", {'color': MUTED, 'size': 13})])

# ---- save ----
out_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "public")
out = os.path.join(out_dir, "강의안_데이터분석_시각화.pptx")
prs.save(out)
print(f"SAVED: {out}  ({n} slides)")
